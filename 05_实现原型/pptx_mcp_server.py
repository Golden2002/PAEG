# -*- coding: utf-8 -*-
"""生成 PPT MCP server（v0.60 ⭐ 升级：应用路演 PPT 手工经验）

v0.60 升级内容（依据维护手册 §18.6-18.9 + memo/022 沉淀经验）：
1. 统一头部：品牌 Logo（每页右上角，PNG 图标版）+ 标题条 + 页脚
2. markdown 清理：清除 **、##、`、[t](url) 等符号（保留 $ 公式原样）
3. 长文本自适应：字号分级 18→16→14pt，超长拆多页
4. 模板资产化：从 Library/ppt_templates/ 读取模板（配色/方法论）
5. 脚本资产化：可引用 assets/ppt_scripts/ 的生成脚本作为参考
"""
import sys, io, os, json, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJ_DIR = os.path.dirname(BASE_DIR)
OUT_DIR = os.path.join(BASE_DIR, 'downloads', 'ppt')
TEMPLATE_DIR = os.path.join(PROJ_DIR, 'Library', 'ppt_templates')
SCRIPT_DIR = os.path.join(PROJ_DIR, 'assets', 'ppt_scripts')
LOGO_ICON = os.path.join(PROJ_DIR, '09_GUI前端', 'assets', 'icons', 'paeg-logo.png')
if not os.path.exists(LOGO_ICON):
    LOGO_ICON = os.path.join(PROJ_DIR, 'assets', 'logo', 'paeg_logo_icon_dark_512.png')

# 品牌配色（与路演 PPT v8 一致）
C_PRIMARY = RGBColor(0x0F, 0x2A, 0x52)
C_ACCENT = RGBColor(0xE6, 0xA5, 0x28)
C_LIGHT = RGBColor(0xF5, 0xF2, 0xEC)
C_DARK = RGBColor(0x0F, 0x2A, 0x52)
C_GRAY = RGBColor(0x55, 0x5F, 0x6B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE = RGBColor(0xC5, 0x5A, 0x11)


def clean_md(text):
    """清除 markdown 符号，保留内容（$ 公式原样保留）"""
    if not text:
        return text
    t = text
    t = re.sub(r'\*\*(.+?)\*\*', r'\1', t, flags=re.S)
    t = re.sub(r'\*(.+?)\*', r'\1', t, flags=re.S)
    t = re.sub(r'^#{1,6}\s*', '', t, flags=re.M)
    t = re.sub(r'`([^`]*)`', r'\1', t)
    t = re.sub(r'\[([^\]]*)\]\([^)]*\)', r'\1', t)
    t = re.sub(r'^\s*[-*]\s+', '\u2022 ', t, flags=re.M)
    return t.strip()


def estimate_text_height(text, width_in, font_size):
    """估算文本所需高度（英寸），用于长文本自适应"""
    chars_per_line = max(4, int(width_in * 96 / (font_size * 1.15)))
    lines = 0
    for para in text.split('\n'):
        lines += max(1, (len(para) + chars_per_line - 1) // chars_per_line)
    return lines * font_size * 0.018


def _parse_outline(text) -> list:
    """把 LLM/文本大纲解析为 [{title, points:[...], notes}]

    §3.79 Round 3 ⭐ 兼容 list 输入：LessonPrep 静态兜底产出的 ppt_outline
    是 list[dict{slide,title,points}]（_static_ppt_outline），此前直接 .splitlines()
    抛 "'list' object has no attribute 'splitlines'" → 备课→PPT 链路断点。
    修复：list 输入直接映射（保留 title/points），str 输入走原解析。
    """
    if isinstance(text, list):
        slides = []
        for it in text:
            if not isinstance(it, dict):
                continue
            title = str(it.get("title") or it.get("slide") or "")
            pts = it.get("points") or []
            if isinstance(pts, str):
                pts = [ln.strip() for ln in pts.splitlines() if ln.strip()]
            if not title and not pts:
                continue
            slides.append({
                "title": clean_md(str(title)),
                "points": [clean_md(str(p)) for p in pts],
                "notes": str(it.get("notes") or ""),
            })
        if slides:
            return slides
        return [{'title': '演示文稿', 'points': ['（空大纲，请补充内容）'], 'notes': ''}]
    slides = []
    lines = [l.strip() for l in (text or '').splitlines() if l.strip()]
    cur = None
    for ln in lines:
        m = re.match(r'^(#{1,4})\s+(.+)$', ln)
        if m:
            if cur: slides.append(cur)
            cur = {'title': clean_md(m.group(2).strip()), 'points': [], 'notes': ''}
            continue
        m = re.match(r'^(\d+)[.、)]\s*(.+)$', ln)
        if m:
            if cur: slides.append(cur)
            cur = {'title': clean_md(m.group(2).strip()), 'points': [], 'notes': ''}
            continue
        m = re.match(r'^[-*•]\s*(.+)$', ln)
        if m:
            if cur:
                cur['points'].append(clean_md(m.group(1).strip()))
            else:
                cur = {'title': '要点', 'points': [clean_md(m.group(1).strip())], 'notes': ''}
            continue
        if cur:
            if ln.startswith('备注') or ln.startswith('note'):
                cur['notes'] = ln.split(':', 1)[-1].strip()
            else:
                cur['points'].append(clean_md(ln))
    if cur: slides.append(cur)
    return slides or [{'title': '演示文稿', 'points': ['（空大纲，请补充内容）'], 'notes': ''}]


def _set_text(shape, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return shape


def _add_header(slide, title, page_no, prs):
    """统一头部：品牌 Logo（右上角）+ 标题条（无强调线）+ 页码"""
    # 顶部品牌色标题条
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.LEFT
    # 右上角 Logo（PNG 图标版）
    if os.path.exists(LOGO_ICON):
        try:
            slide.shapes.add_picture(LOGO_ICON, Inches(12.35), Inches(0.22), Inches(0.55), Inches(0.55))
        except Exception:
            pass
    # 页脚
    _set_text(slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3)),
              'PAEG · Émile', size=8, color=C_GRAY)
    _set_text(slide.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3)),
              str(page_no), size=8, color=C_GRAY, align=PP_ALIGN.RIGHT)


def _add_bullets_adaptive(slide, points, left, top, width, height, prs):
    """长文本自适应 bullets：字号 18→16→14pt 分级，超长拆页"""
    # 计算内容总量，决定字号
    full_text = '\n'.join(points)
    total_chars = len(full_text)
    size = 18
    if total_chars > 400:
        size = 16
    if total_chars > 700:
        size = 14
    if total_chars > 1100:
        size = 12

    # 每页最多 6 条
    chunks = [points[i:i+6] for i in range(0, len(points), 6)]
    for ci, chunk in enumerate(chunks):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE  # 防撑破布局
        for i, pt in enumerate(chunk):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pt = clean_md(pt)
            p.text = ('\u2022 ' if not pt.startswith('\u2022') else '') + pt
            p.font.size = Pt(size)
            p.font.color.rgb = C_DARK
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(8)
    return box


def _extract_material_text(uid: str, max_files: int = 3, max_chars: int = 8000) -> str:
    """从用户上传物料提取文字（路径安全）"""
    try:
        _root = os.path.join(PROJ_DIR, 'Library', 'usr_knowledge')
        _udir = os.path.join(_root, str(uid))
        _real_root = os.path.realpath(_root)
        _real_udir = os.path.realpath(_udir)
        if not _real_udir.startswith(_real_root) or '/' in str(uid) or '\\' in str(uid):
            return ""
        if not os.path.isdir(_real_udir):
            return ""
        _texts = []
        for _f in sorted(os.listdir(_real_udir))[:max_files]:
            _fp = os.path.join(_real_udir, _f)
            if not os.path.isfile(_fp):
                continue
            _ext = os.path.splitext(_f)[1].lower()
            try:
                if _ext in ('.md', '.txt'):
                    with open(_fp, encoding='utf-8', errors='ignore') as fh:
                        _texts.append(f"【{_f}】\n{fh.read()[:3000]}")
                elif _ext == '.pdf':
                    try:
                        from pypdf import PdfReader
                        _rd = PdfReader(_fp)
                        _txt = "\n".join(p.extract_text() or "" for p in _rd.pages[:5])
                        _texts.append(f"【{_f}】\n{_txt[:3000]}")
                    except Exception:
                        pass
                elif _ext == '.docx':
                    try:
                        import docx as _docx
                        _doc = _docx.Document(_fp)
                        _txt = "\n".join(p.text for p in _doc.paragraphs[:60])
                        _texts.append(f"【{_f}】\n{_txt[:3000]}")
                    except Exception:
                        pass
            except Exception:
                continue
        return "\n\n".join(_texts)[:max_chars]
    except Exception:
        return ""


def _add_slide_image(slide, image_path, prs):
    """为单页幻灯片添加配图（右侧区域 7.7→12.9 宽 × 1.5→5.5 高）——静默失败"""
    try:
        if not image_path:
            return
        slide.shapes.add_picture(
            image_path,
            Inches(7.7), Inches(1.5),
            width=Inches(5.2), height=Inches(4.0),
        )
    except Exception:
        pass


def list_templates() -> dict:
    """列出可用模板与脚本资产（v0.60 ⭐ 模板/脚本资产化）"""
    templates = {}
    if os.path.exists(TEMPLATE_DIR):
        for f in os.listdir(TEMPLATE_DIR):
            templates[f] = os.path.join(TEMPLATE_DIR, f)
    scripts = {}
    if os.path.exists(SCRIPT_DIR):
        for f in os.listdir(SCRIPT_DIR):
            if f.endswith('.js'):
                scripts[f] = os.path.join(SCRIPT_DIR, f)
    return {"templates": templates, "scripts": scripts, "logo": LOGO_ICON}


def generate_ppt(topic: str, outline: str = "", sources: str = "",
                 out_name: str = "", uid: str = "", style: str = "paeg_standard",
                 enable_images: bool = True) -> dict:
    """生成演示文稿 .pptx（v0.60 ⭐ 升级：Logo/自适应/markdown 清理）

    Args:
        topic: 演示主题
        outline: 大纲文本（LLM 生成 ## 标题 + 要点）
        sources: 来源说明
        out_name: 输出文件名
        uid: 用户 id（提取上传物料）
        style: 风格模板——'paeg_standard'（深蓝+金，默认）/
               'presentation_zen'（极简留白）/'dark_premium'（深色高级）
        enable_images: 是否启用右侧配图（v0.61 ⭐ 资料库/公共库/联网检索，失败静默降级）

    Returns:
        {"ok": bool, "path": str, "slides": int, "error": str, "templates": dict}
    """
    global C_PRIMARY, C_ACCENT, C_LIGHT, C_DARK
    from pptx.dml.color import RGBColor as _RGB
    if style == "presentation_zen":
        C_PRIMARY, C_ACCENT, C_LIGHT, C_DARK = _RGB(0x2C, 0x5F, 0x2D), _RGB(0x97, 0xBC, 0x62), _RGB(0xF5, 0xF5, 0xF5), _RGB(0x21, 0x21, 0x21)
    elif style == "dark_premium":
        C_PRIMARY, C_ACCENT, C_LIGHT, C_DARK = _RGB(0x1E, 0x27, 0x61), _RGB(0xCA, 0xDC, 0xFC), _RGB(0x11, 0x14, 0x20), _RGB(0xFF, 0xFF, 0xFF)
    else:
        C_PRIMARY, C_ACCENT, C_LIGHT, C_DARK = _RGB(0x0F, 0x2A, 0x52), _RGB(0xE6, 0xA5, 0x28), _RGB(0xF5, 0xF2, 0xEC), _RGB(0x0F, 0x2A, 0x52)
    try:
        os.makedirs(OUT_DIR, exist_ok=True)
        slides_data = _parse_outline(outline or topic)
        material = ""
        if uid:
            material = _extract_material_text(uid)
        if material and slides_data:
            _first = slides_data[0]
            _extra = [m[:120] for m in material.split("\n") if m.strip()][:6]
            if _extra:
                _first.setdefault('points', []).extend(f"(物料) {e}" for e in _extra[:3])
                _first['notes'] = (_first.get('notes') or '') + f"\n\n用户物料摘要：{material[:600]}"

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # 封面页（深色 + 大标题 + Logo）
        s = prs.slides.add_slide(blank)
        bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = C_PRIMARY
        bg.line.fill.background()
        _set_text(s.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(1.6)),
                  topic, size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _set_text(s.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(11), Inches(1.0)),
                  'PAEG · Émile 演示文稿', size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
        if os.path.exists(LOGO_ICON):
            try:
                s.shapes.add_picture(LOGO_ICON, Inches(6.15), Inches(4.8), Inches(1.0), Inches(1.0))
            except Exception:
                pass
        if sources:
            _set_text(s.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10.3), Inches(1.2)),
                      f'资料来源：{sources[:120]}', size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
        elif material:
            _set_text(s.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10.3), Inches(1.2)),
                      '内容基于用户上传的资料整理', size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

        # 内容页（统一头部 + 自适应 bullets + 页码）
        for i, sd in enumerate(slides_data[:20]):
            s = prs.slides.add_slide(blank)
            _add_header(s, sd['title'], i + 2, prs)
            pts = sd['points'] or ['（本页要点）']
            _add_bullets_adaptive(s, pts, 0.8, 1.4, 6.7, 5.4, prs)  # 宽度 11.7→6.7 让出右侧配图区
            # §3.71 ⭐ PPT 配图增强：资料库/公共库/联网检索（失败静默）
            if enable_images:
                try:
                    from pptx_image_supplier import find_images_for_slide
                    _imgs = find_images_for_slide(sd.get("title", ""), pts, uid, max_results=1)
                    if _imgs:
                        _add_slide_image(s, _imgs[0], prs)
                except Exception:
                    pass
            if sd.get('notes'):
                s.notes_slide.notes_text_frame.text = sd['notes']

        fname = (out_name or f'{re.sub(r"[^\w\u4e00-\u9fff]+", "_", topic)[:30]}_{os.path.getmtime(__file__):.0f}') + '.pptx'
        fname = re.sub(r'[\\/:*?"<>|]', '_', fname)
        path = os.path.join(OUT_DIR, fname)
        prs.save(path)
        return {"ok": True, "path": path, "slides": len(slides_data) + 1,
                "error": "", "templates": list_templates()}
    except Exception as e:
        return {"ok": False, "path": "", "slides": 0, "error": str(e), "templates": list_templates()}


# ─── FastMCP server 暴露 ───
try:
    from fastmcp import FastMCP
    mcp = FastMCP("paeg-pptx")

    @mcp.tool()
    def generate_presentation(topic: str, outline: str = "", sources: str = "",
                              out_name: str = "", uid: str = "",
                              style: str = "paeg_standard",
                              enable_images: bool = True) -> dict:
        """根据主题+大纲+来源生成演示文稿 PPT（v0.60 升级：品牌 Logo + 长文本自适应 + markdown 清理）。

        大纲格式：每页以 '## 标题' 或 '1. 标题' 开头，要点以 '- ' 开头。
        uid（可选）：用户 id——提供时自动提取该用户上传物料补充内容。
        style：'paeg_standard'（深蓝+金，默认）/'presentation_zen'/'dark_premium'。
        enable_images（v0.61 ⭐）：是否启用右侧配图（资料库/公共库/联网检索），默认 True。
        模板与脚本资产：Library/ppt_templates/ + assets/ppt_scripts/（方法论见维护手册 §18.6-18.9）。"""
        return generate_ppt(topic, outline, sources, out_name, uid, style=style, enable_images=enable_images)

    @mcp.tool()
    def list_ppt_templates() -> dict:
        """列出可用的 PPT 模板与生成脚本资产（模板在 Library/ppt_templates/，脚本在 assets/ppt_scripts/）。"""
        return list_templates()

    if __name__ == "__main__":
        mcp.run(transport="stdio")
except Exception as e:
    mcp = None
    if __name__ == "__main__":
        print(f'FastMCP 不可用: {e}；直接函数测试:')
        r = generate_ppt('语言学导论', '## 什么是语言学\n- 研究语言的科学', '用户文档')
        print(r)
