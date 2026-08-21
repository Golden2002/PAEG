"""
PAEG 知识库扩展加载器（v0.11）

任务 C：为智能体留出"扩展接口 + 文件夹"，允许未来向 Library 中加入大量知识库，
为模型提供更多现实的事实支撑。

目录约定（Library/ 下）：
  Library/KnowledgeBase/subjects/  *.json  —— 学科知识节点（与 knowledge_base.py 同构）
  Library/KnowledgeBase/facts/     *.md    —— 事实性资料（现实数据/史实/原文摘录）
  Library/Language/   *.md         —— 词汇/语法（未来可索引）
  Library/Math/       *.pdf        —— 数学资料（未来可解析）
  Library/Philosophy/ *.pdf        —— 哲学文本（未来可解析）
  Library/Simone Weil/ *.docx/pdf  —— 薇依原文（未来可检索）

用法：
    from library_loader import KnowledgeLibrary
    kl = KnowledgeLibrary()          # 扫描并加载
    kl.register()                    # 把学科节点并入 KnowledgeBase
    kl.search_facts("光合作用")       # 检索事实资料
"""

from __future__ import annotations

import glob
import json
import os
from typing import Dict, List, Optional


class KnowledgeLibrary:
    """扫描 Library 目录，加载可用的知识源。

    §3.79 Round 11 ⭐ 性能优化：PDF/PPTX 课件文本提取较重（22 份 PDF 全量提取
    耗时数十秒）——模块级单例缓存（_get_instance），服务器/测试进程内只提取一次。
    """

    _instance = None

    @classmethod
    def _get_instance(cls, base_dir=None):
        if cls._instance is None:
            cls._instance = cls(base_dir=base_dir)
        return cls._instance

    def __init__(self, base_dir: Optional[str] = None):
        # 默认 Library 目录
        self.base_dir = base_dir or os.path.join(
            r'D:\桌面\智能体架构与开发（含大模型）\14_教育者Agent项目', 'Library')
        self.subjects: Dict[str, dict] = {}   # 从 subjects/*.json 加载
        self.facts: List[dict] = {}           # 从 facts/*.md 加载（文件名->内容）
        self.raw_files: List[dict] = []       # 其他可索引的源文件
        self._scan()

    def _scan(self):
        """扫描目录结构。"""
        kb_dir = os.path.join(self.base_dir, 'KnowledgeBase')
        # 1. 学科节点 JSON
        for fp in glob.glob(os.path.join(kb_dir, 'subjects', '*.json')):
            try:
                with open(fp, encoding='utf-8') as f:
                    data = json.load(f)
                if isinstance(data, dict):
                    for nid, node in data.items():
                        if isinstance(node, dict):
                            node.setdefault('id', nid)
                            self.subjects[nid] = node
            except Exception as e:
                print(f'[library] 加载学科节点失败 {fp}: {e}')
        # 2. 事实资料 MD（文件名作为标题）
        for fp in glob.glob(os.path.join(kb_dir, 'facts', '*.md')):
            name = os.path.splitext(os.path.basename(fp))[0]
            try:
                with open(fp, encoding='utf-8') as f:
                    content = f.read()
                self.facts[name] = content
            except Exception as e:
                print(f'[library] 加载事实资料失败 {fp}: {e}')
        # 3. 其他源文件（仅登记路径，未来接入解析）
        # §3.79 Round 11 ⭐ 增强：.pptx/.pdf 课件文本提取——张宇扬课件知识库
        # （Library/common/张宇扬课件/）主体为 .pdf/.pptx（466MB 26 文件），此前仅登记
        # 不解析 → 内容不可检索。现用 python-pptx / pypdf 提取并入 raw_files（content 字段）。
        # ⭐ 落盘缓存：提取较重（34 份 PDF 全量数十秒）——结果按 文件mtime 缓存到
        # Library/.courseware_cache.json，未变更则跳过提取（首次提取后秒级加载）。
        # v1.2.23 ⭐ 性能再优化：缓存含 _manifest（相对路径→mtime），二次扫描直接按
        # 清单构建（零 os.walk / 零 stat）——慢速磁盘上对 466MB 大文件逐个 getmtime
        # 实测需 36-83s（杀软/同步干扰），清单命中后构造降至毫秒级。
        _cache_path = os.path.join(self.base_dir, '.courseware_cache.json')
        _cache = {}
        try:
            if os.path.isfile(_cache_path):
                with open(_cache_path, encoding='utf-8') as _cf:
                    _cache = json.load(_cf)
        except Exception:
            _cache = {}
        _need_save = False
        _manifest = _cache.get('_manifest') or {}

        def _cached_entry(fp: str, f: str, root: str, kind: str, extractor) -> dict:
            """带 mtime 缓存的课件提取（命中缓存跳过 extractor）。"""
            nonlocal _need_save
            _entry = {'path': fp, 'name': f, 'category': os.path.basename(root)}
            try:
                _mtime = os.path.getmtime(fp)
            except Exception:
                _mtime = 0
            _key = os.path.relpath(fp, self.base_dir)
            _cached = _cache.get(_key)
            if _cached and abs(float(_cached.get('mtime', 0)) - _mtime) < 1:
                if _cached.get('content'):
                    _entry['content'] = _cached['content']
                return _entry
            try:
                _content = extractor()
                if _content:
                    _entry['content'] = _content
                    _cache[_key] = {'mtime': _mtime, 'content': _content,
                                    'kind': kind}
                    _need_save = True
            except Exception as _ex_e:
                print(f"[library] {kind} 提取跳过 {f}: {_ex_e}")
            return _entry

        # ── 快速路径：缓存清单命中 → 直接按清单构建（零 walk/stat）──
        if _manifest:
            for _rel, _info in _manifest.items():
                try:
                    _fp = os.path.join(self.base_dir, _rel)
                    if not os.path.isfile(_fp):
                        continue
                except Exception:
                    continue
                _f = os.path.basename(_fp)
                _cat = os.path.basename(os.path.dirname(_fp))
                _entry = {'path': _fp, 'name': _f, 'category': _cat}
                _cached = _cache.get(_rel)
                if _cached and _cached.get('content'):
                    _entry['content'] = _cached['content']
                self.raw_files.append(_entry)
            return  # 快速路径结束（无需 walk）

        for root, dirs, files in os.walk(self.base_dir):
            if 'KnowledgeBase' in root:
                continue
            for f in files:
                if f.endswith(('.md', '.txt', '.json')):
                    fp = os.path.join(root, f)
                    self.raw_files.append({'path': fp, 'name': f,
                                           'category': os.path.basename(root)})
                elif f.endswith('.pptx'):
                    fp = os.path.join(root, f)

                    def _extract_pptx(_fp=fp):
                        from pptx import Presentation
                        _prs = Presentation(_fp)
                        _pages = []
                        for _slide in _prs.slides:
                            _texts = []
                            for _shape in _slide.shapes:
                                if getattr(_shape, 'has_text_frame', False):
                                    _t = _shape.text_frame.text.strip()
                                    if _t:
                                        _texts.append(_t)
                                elif getattr(_shape, 'has_table', False):
                                    for _row in _shape.table.rows:
                                        for _cell in _row.cells:
                                            _ct = _cell.text.strip()
                                            if _ct:
                                                _texts.append(_ct)
                            if _texts:
                                _pages.append(" | ".join(_texts))
                        if _pages:
                            print(f"[library] pptx 课件已提取: {os.path.basename(_fp)} "
                                  f"({len(_pages)} 页)")
                            return "\n".join(_pages)
                        return ""

                    self.raw_files.append(
                        _cached_entry(fp, f, root, "pptx", _extract_pptx))
                elif f.endswith('.pdf'):
                    # §3.79 Round 11 ⭐ PDF 课件文本提取（张宇扬课件主体为 .pdf 22 份）：
                    # 用 pypdf 逐页提取文本；扫描版/加密 PDF 提取失败 → 仅登记（不阻塞）。
                    fp = os.path.join(root, f)

                    def _extract_pdf(_fp=fp):
                        from pypdf import PdfReader
                        _reader = PdfReader(_fp)
                        _pages = []
                        for _p in _reader.pages:
                            try:
                                _t = (_p.extract_text() or "").strip()
                            except Exception:
                                _t = ""
                            if _t:
                                _pages.append(_t)
                        if _pages:
                            print(f"[library] pdf 课件已提取: {os.path.basename(_fp)} "
                                  f"({len(_pages)} 页)")
                            return "\n".join(_pages)
                        return ""

                    self.raw_files.append(
                        _cached_entry(fp, f, root, "pdf", _extract_pdf))
        if _need_save or not _manifest:
            # 写回时记录文件清单（快速路径依赖）
            for _entry in self.raw_files:
                _rel = os.path.relpath(_entry['path'], self.base_dir)
                if _rel not in _manifest:
                    try:
                        _manifest[_rel] = {'mtime': os.path.getmtime(_entry['path']),
                                           'kind': 'file'}
                    except Exception:
                        pass
            _cache['_manifest'] = _manifest
            try:
                with open(_cache_path, 'w', encoding='utf-8') as _cf:
                    json.dump(_cache, _cf, ensure_ascii=False)
            except Exception as _cse:
                print(f"[library] 课件缓存写入失败: {_cse}")

    def register(self, knowledge_base) -> int:
        """把加载的学科节点并入 PAEG 的 KnowledgeBase。返回新增节点数。"""
        added = 0
        for nid, node in self.subjects.items():
            if nid not in knowledge_base.subjects:
                knowledge_base.subjects[nid] = node
                added += 1
        return added

    def search_facts(self, query: str, top_k: int = 3) -> List[dict]:
        """在事实资料中检索关键词。返回匹配片段。

        §3.79 Round 11 ⭐ 增强：除 facts/*.md 外，也检索 raw_files 中带 content 的
        课件文本（pptx 提取页/README 等）——张宇扬课件知识库真正可检索。
        """
        q = query.lower()
        results = []
        for name, content in self.facts.items():
            if q in name.lower() or q in content.lower():
                # 找匹配段落
                for para in content.split('\n\n'):
                    if q in para.lower():
                        results.append({'source': name, 'snippet': para[:300]})
                        break
        # Round 11 ⭐ 课件文本检索（raw_files.content）
        if len(results) < top_k:
            for entry in self.raw_files:
                _content = entry.get("content") or ""
                if not _content:
                    continue
                _name = str(entry.get("name") or "")
                _cat = str(entry.get("category") or "")
                if q not in _content.lower() and q not in _name.lower():
                    continue
                # 找匹配行（课件按页换行）
                _hit = None
                for _pg in _content.split("\n"):
                    if q in _pg.lower():
                        _hit = _pg[:300]
                        break
                results.append({
                    "source": f"{_cat}/{_name}",
                    "snippet": _hit or _content[:200],
                    "kind": "courseware",
                })
                if len(results) >= top_k:
                    break
        return results[:top_k]

    def list_sources(self) -> List[dict]:
        """列出所有可用的知识源。"""
        return [
            {'type': 'subject', 'name': nid} for nid in self.subjects
        ] + [
            {'type': 'fact', 'name': name} for name in self.facts
        ] + self.raw_files

    def stats(self) -> dict:
        return {
            'subjects': len(self.subjects),
            'facts': len(self.facts),
            'raw_files': len(self.raw_files),
        }
