# -*- coding: utf-8 -*-
"""资源加载器（v1.0 ⭐）：Library 用户物料 + PPT 公共模板主色。

安全约束：load_user_asset 做 realpath 路径校验（防目录穿越），
只允许读取 <Library/usr_knowledge/<uid>>/ 内的文件。
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Dict, Optional, List

_PROJ = Path(__file__).resolve().parents[1]
_USR_ROOT = _PROJ / "Library" / "usr_knowledge"
_TEMPLATE_DIR = _PROJ / "Library" / "ppt_templates"

_TEMPLATE_COLORS: Dict[str, dict] = {}


def load_user_asset(uid: str, path: str) -> Optional[str]:
    """路径安全解析：相对路径必须落在 <USR_ROOT>/<uid>/ 下，返回绝对路径。

    穿越/不存在 → None（调用方降级处理，不阻断视频）。
    """
    if not uid or not path:
        return None
    try:
        base = (_USR_ROOT / str(uid)).resolve()
        target = (base / path).resolve()
    except Exception:
        return None
    if not str(target).startswith(str(base)):
        return None
    return str(target) if target.is_file() else None


def list_user_assets(uid: str) -> List[str]:
    """列出该用户 Library/usr_knowledge/<uid>/ 下全部文件。"""
    d = _USR_ROOT / str(uid)
    if not d.is_dir():
        return []
    return [str(p) for p in d.rglob("*") if p.is_file()]


def get_template_colors(style: str) -> dict:
    """从 Library/ppt_templates/<style>.pptx 读主色（python-pptx，缓存）。

    返回 {"primary": "RRGGBB", ...}；解析失败返回 {}（帧绘制用默认色）。
    """
    if style in _TEMPLATE_COLORS:
        return _TEMPLATE_COLORS[style]
    p = _TEMPLATE_DIR / f"{style}.pptx"
    cache: dict = {}
    if p.is_file():
        try:
            from pptx import Presentation
            prs = Presentation(str(p))
            if prs.slides:
                shp = prs.slides[0].shapes[0]
                fill = getattr(shp, "fill", None)
                if fill is not None and fill.type == 1:  # solid
                    rgb = fill.fore_color.rgb
                    cache = {"primary": str(rgb)}
        except Exception:
            cache = {}
    _TEMPLATE_COLORS[style] = cache
    return cache
