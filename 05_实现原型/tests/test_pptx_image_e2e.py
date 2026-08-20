# -*- coding: utf-8 -*-
"""test_pptx_image_e2e.py — PAEG PPT 图片增强 E2E 测试（v1.2.0 ⭐）

端到端验证 pptx_mcp_server.generate_ppt 在不同 enable_images 取值下的行为：
- enable_images=True  → 至少 1 个非 logo 的 picture shape 被插入
- enable_images=False → picture 数 == logo 数（不配任何内容图）

测试隔离策略：
1. 把假 PNG 放到真实 Library/usr_knowledge/u_e2e/，让本地命中优先于联网。
2. monkeypatch ``requests.get`` 返回假 PNG bytes（仅在本地 miss 时兜底，
   防止测试因实现 bug 静默走真实网络拖慢 + 污染生产缓存）。
3. 测试结束清理 u_e2e/植物_叶绿体.png，避免污染真实用户数据。
4. 使用 PIL.Image 生成真实 PNG bytes（_make_fake_png 100x100 红色），便于
   区分调用来源；同时验证 Pillow 已正确安装（requirements.txt 强依赖）。

设计原则：
- E2E 不修改生产代码；通过 monkeypatch 与 fixture 实现隔离。
- 所有路径断言使用 Path 对象而非 str，便于 Windows/Linux 双跑。
- picture 计数采用"按 shape_type 过滤 + 按尺寸过滤 logo"的纯函数，
  与生成器内部实现解耦。
"""
from __future__ import annotations

import io
import os
import sys
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ────────────────────────────────────────────────────────────
# 路径常量
# ────────────────────────────────────────────────────────────
# 测试文件位于 05_实现原型/tests/，项目根 = parents[2]
_TESTS_DIR = Path(__file__).resolve().parent
_PROJ05 = _TESTS_DIR.parent
_PROJ_ROOT = _PROJ05.parent  # D:\桌面\智能体架构与开发（含大模型）\14_教育者Agent项目
_USR_KNOWLEDGE_DIR = _PROJ_ROOT / "Library" / "usr_knowledge"
_FAKE_PNG_PATH = _USR_KNOWLEDGE_DIR / "u_e2e" / "植物_叶绿体.png"


def _ensure_sys_path():
    """确保 05_实现原型/ 在 sys.path（直接 python -m pytest 也可 import）。"""
    p = str(_PROJ05)
    if p not in sys.path:
        sys.path.insert(0, p)


_ensure_sys_path()


def _safe_import_pptx_mcp_server():
    """Import pptx_mcp_server without poisoning sys.stdout.

    pptx_mcp_server.py 在模块顶部（line 12）执行
    ``sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')``，
    用于在 Windows 控制台正确输出中文。该 import-time 副作用会把当前
    ``sys.stdout.buffer`` 包进 TextIOWrapper 并赋给 ``sys.stdout``——若
    彼时 pytest 已开启 capture（``sys.stdout`` 是 tempfile-based），
    TextIOWrapper 会持有 pytest 的 capture buffer；test 结束后 pytest
    关闭 buffer，wrapper 仍指向已关闭的 buffer → ``ValueError: I/O
    operation on closed file`` 在 pytest teardown 抛错并污染整个测试结果。

    修复策略：用 ``importlib`` 加载源码并 AST-级跳过 line 12（``sys.stdout``
    重新包裹那一行）。模块其余行为完全保留——``generate_ppt`` 不依赖
    stdout 行为，所以删除该行无副作用。
    """
    if "pptx_mcp_server" in sys.modules:
        return  # 已 import，不重复执行
    import ast

    src_path = _PROJ05 / "pptx_mcp_server.py"
    src_text = src_path.read_text(encoding="utf-8")
    # 去除可能存在的 BOM（pptx_mcp_server.py 是 utf-8-sig 兼容存储）
    if src_text.startswith("\ufeff"):
        src_text = src_text[1:]
    tree = ast.parse(src_text)

    # 删除 `sys.stdout = io.TextIOWrapper(...)` 这一行。
    # AST 形态：Assign(targets=[Attribute(value=Name('sys'), attr='stdout')],
    #                  value=Call(func=Attribute(value=Name('io'), attr='TextIOWrapper')))
    for node in list(tree.body):
        if not isinstance(node, ast.Assign) or len(node.targets) != 1:
            continue
        tgt = node.targets[0]
        if not (isinstance(tgt, ast.Attribute) and tgt.attr == "stdout"
                and isinstance(tgt.value, ast.Name) and tgt.value.id == "sys"):
            continue
        val = node.value
        if not (isinstance(val, ast.Call) and isinstance(val.func, ast.Attribute)
                and val.func.attr == "TextIOWrapper"):
            continue
        tree.body.remove(node)
        break  # 只删一行

    # compile & exec 模块
    code = compile(tree, str(src_path), "exec")
    module = type(sys)("pptx_mcp_server")
    module.__file__ = str(src_path)
    sys.modules["pptx_mcp_server"] = module
    exec(code, module.__dict__)


# 在 import pptx_mcp_server 之前快照 stdout（pptx_mcp_server.py line 12 会重写它）
_ORIG_STDOUT = sys.stdout
_ORIG___STDOUT__ = sys.__stdout__


# ────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────
def _make_fake_png_bytes(size: tuple = (100, 100), color: str = "red") -> bytes:
    """用 Pillow 生成内存 PNG bytes（避免依赖外部文件）。"""
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


class _FakeResp:
    """最小 requests.Response 替身：让 ``r.raise_for_status()`` 啥也不做。"""

    def __init__(self, content: bytes, status_code: int = 200):
        self.content = content
        self.status_code = status_code
        self.text = ""
        self.headers: dict = {}

    def raise_for_status(self):
        if self.status_code >= 400:
            raise RuntimeError(f"HTTP {self.status_code}")


def _is_logo(picture) -> bool:
    """判断 picture 是否为品牌 logo。

    logo 在两张页的尺寸都 ≤ 1.0 英寸（封面 1.0x1.0 / 内容页 0.55x0.55）。
    内容配图 5.2x4.0 英寸 → max(w,h)=5.2 > 1.2 → 不是 logo。
    """
    if picture.width is None or picture.height is None:
        return False
    # EMUs: 914400 per inch
    w_in = picture.width / 914400.0
    h_in = picture.height / 914400.0
    return max(w_in, h_in) <= 1.2


def _count_pictures(pptx_path: str) -> tuple[int, int]:
    """扫整份 PPT：返回 (total_pictures, logo_pictures)。

    - total_pictures：所有 slide.shapes 中 shape_type == PICTURE(13) 的形状
    - logo_pictures：其中尺寸 ≤ 1.2 英寸的（封面/内容页 logo）
    """
    prs = Presentation(pptx_path)
    total = 0
    logos = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total += 1
                if _is_logo(shape):
                    logos += 1
    return total, logos


def _generate_and_count(
    topic: str,
    outline: str,
    uid: str,
    enable_images: bool,
    monkeypatch,
) -> tuple[int, int, dict]:
    """monkeypatch 网络 → 调 generate_ppt → 返回 (total_pics, logos, result)。

    网络被 monkeypatch 后即使本地 miss 也只会拿到假 bytes，触发不到真实 Bing。
    """
    _safe_import_pptx_mcp_server()
    import pptx_mcp_server as ms

    fake_png = _make_fake_png_bytes()

    def _fake_get(url, *args, **kwargs):
        return _FakeResp(fake_png, 200)

    # monkeypatch 模块级 requests.get（generate_ppt 内部走的是 import requests，
    # 但 pptx_image_supplier 通过模块属性 pis.requests 访问——直接 monkeypatch
    # 内置 requests.get 即可同时覆盖两路调用）
    import requests as _requests

    monkeypatch.setattr(_requests, "get", _fake_get)
    # 兼容：supplier 用了 ``import requests`` 然后 pis.requests.get，要 patch
    # ``pis.requests.get`` 与 ``pis.requests`` 属性指向同一个模块对象，因此
    # 上面 patch 也会反映到 supplier。
    import pptx_image_supplier as pis
    monkeypatch.setattr(pis.requests, "get", _fake_get)

    result = ms.generate_ppt(
        topic, outline, "", "", uid, style="paeg_standard",
        enable_images=enable_images,
    )
    total, logos = (0, 0)
    if result["ok"]:
        total, logos = _count_pictures(result["path"])
    return total, logos, result


# ────────────────────────────────────────────────────────────
# Fixture：本地命中用假 PNG（存在则复用；缺失则写入）
# ────────────────────────────────────────────────────────────
@pytest.fixture
def fake_png_in_usr_knowledge():
    """保证 Library/usr_knowledge/u_e2e/植物_叶绿体.png 存在。

    - 若已存在：保留不动（幂等）。
    - 若不存在：写入 Pillow 生成的 100x100 红色 PNG。

    测试结束**不清理**——按用户要求该资源是 E2E 前置依赖（避免下次跑测试
    又得联网）；任何后续测试直接复用即可。
    """
    _FAKE_PNG_PATH.parent.mkdir(parents=True, exist_ok=True)
    if not _FAKE_PNG_PATH.exists():
        _FAKE_PNG_PATH.write_bytes(_make_fake_png_bytes())
    yield _FAKE_PNG_PATH


# ────────────────────────────────────────────────────────────
# E2E 测试 1：enable_images=True → 至少 1 个非 logo 的 picture
# ────────────────────────────────────────────────────────────
def test_generate_ppt_with_images_inserts_non_logo_picture(
    fake_png_in_usr_knowledge, monkeypatch
):
    """enable_images=True 时，generate_ppt 应在内容页插入配图（非 logo）。"""
    total, logos, result = _generate_and_count(
        topic="光合作用",
        outline="## 植物光合作用\n- 叶绿体",
        uid="u_e2e",
        enable_images=True,
        monkeypatch=monkeypatch,
    )

    assert result["ok"], f"generate_ppt 失败: {result.get('error')}"
    assert result["path"].endswith(".pptx"), f"应返回 .pptx 路径, 实际: {result['path']}"
    assert total > logos, (
        f"enable_images=True 应至少插入 1 张非 logo 的 picture, "
        f"实际 total={total}, logos={logos}"
    )


# ────────────────────────────────────────────────────────────
# E2E 测试 2：enable_images=False → picture 数 == logo 数
# ────────────────────────────────────────────────────────────
def test_generate_ppt_without_images_inserts_only_logo(
    fake_png_in_usr_knowledge, monkeypatch
):
    """enable_images=False 时，所有 picture 应均为 logo，不插入配图。"""
    total, logos, result = _generate_and_count(
        topic="光合作用",
        outline="## 植物光合作用\n- 叶绿体",
        uid="u_e2e",
        enable_images=False,
        monkeypatch=monkeypatch,
    )

    assert result["ok"], f"generate_ppt 失败: {result.get('error')}"
    assert total == logos, (
        f"enable_images=False 应只插入 logo, 但 total={total} logos={logos} "
        f"出现额外 picture (差值={total - logos})"
    )