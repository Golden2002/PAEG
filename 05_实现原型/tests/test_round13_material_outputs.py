# -*- coding: utf-8 -*-
"""§3.79 Round 3 ⭐ 物料产出链路修复回归测试。

probe_material_outputs 真实抽查暴露两个断点：
  1. _static_script 讲稿模板缺生活化例子（has_example=False → material_quality 检查失败）
  2. pptx_mcp_server._parse_outline 只收 str，LessonPrep 静态兜底产出 list[dict]
     → "'list' object has no attribute 'splitlines'" → 备课→PPT 链路断

本测试守卫：
  - 静态讲稿模板含生活化例子（开场/主体/小结都带类比）
  - _parse_outline 兼容 list 输入（映射 title/points/notes）
  - list 输入生成真实 .pptx 且可被 python-pptx 打开
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest


def _import_pptx_mcp_server():
    """安全导入 pptx_mcp_server：AST 级跳过其顶层 sys.stdout 重写行。

    与 test_pptx_image_e2e._safe_import_pptx_mcp_server 同款防护——
    该模块 import 时执行 `sys.stdout = io.TextIOWrapper(...)`，
    会污染 pytest capture（teardown 抛 "I/O operation on closed file"）。
    """
    if "pptx_mcp_server" in sys.modules:
        return sys.modules["pptx_mcp_server"]
    import ast

    src_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                            "pptx_mcp_server.py")
    src_text = open(src_path, encoding="utf-8").read()
    if src_text.startswith("\ufeff"):
        src_text = src_text[1:]
    tree = ast.parse(src_text)
    for node in list(tree.body):
        if not isinstance(node, ast.Assign) or len(node.targets) != 1:
            continue
        tgt = node.targets[0]
        if not (isinstance(tgt, ast.Attribute) and tgt.attr == "stdout"
                and isinstance(tgt.value, ast.Name) and tgt.value.id == "sys"):
            continue
        val = node.value
        if not (isinstance(val, ast.Call) and isinstance(val.func, ast.Attribute)
                and val.func.attr == "TextIOWrapper"):
            continue
        tree.body.remove(node)
        break
    code = compile(tree, str(src_path), "exec")
    module = type(sys)("pptx_mcp_server")
    module.__file__ = str(src_path)
    sys.modules["pptx_mcp_server"] = module
    exec(code, module.__dict__)
    return module


def test_static_script_has_life_example():
    """静态讲稿模板必须含生活化例子（Round 3 修复：check_lecture_script.has_example）。"""
    from subagents import _static_script
    from services.material_quality import check_lecture_script

    script = _static_script("函数的单调性", 15)
    chk = check_lecture_script(script)
    assert chk["passed"], f"静态讲稿应过检查器: {chk['errors']}"
    assert chk["has_example"], "静态讲稿缺生活化例子（has_example=False）"
    # 例子的具体形态：类比/例子/比作/就像 至少其一
    assert any(k in script for k in ("例子", "类比", "就像", "比作")), \
        "讲稿中应含 例子/类比/就像/比作 一类生活化表达"


def test_parse_outline_accepts_list():
    """_parse_outline 必须兼容 LessonPrep 的 list 输出（Round 3 修复）。"""
    pms = _import_pptx_mcp_server()

    outline = [
        {"slide": 1, "title": "导入：单调性", "points": ["生活实例", "问题引入"]},
        {"slide": 2, "title": "单调性定义", "points": ["核心概念"], "notes": "强调严格定义"},
        {"slide": 3, "title": "判定方法", "points": ["定义法", "导数法"]},
    ]
    slides = pms._parse_outline(outline)
    assert isinstance(slides, list) and len(slides) == 3
    assert slides[0]["title"] == "导入：单调性"
    assert slides[0]["points"] == ["生活实例", "问题引入"]
    assert slides[1]["notes"] == "强调严格定义"


def test_parse_outline_list_str_points():
    """points 为字符串（非 list）时也应兜底拆行。"""
    pms = _import_pptx_mcp_server()

    outline = [{"slide": 1, "title": "测试", "points": "第一点\n第二点"}]
    slides = pms._parse_outline(outline)
    assert slides[0]["points"] == ["第一点", "第二点"]


def test_parse_outline_empty_list_fallback():
    """空 list → 返回占位单页（不崩溃）。"""
    pms = _import_pptx_mcp_server()

    slides = pms._parse_outline([])
    assert len(slides) == 1
    assert slides[0]["title"]


def test_parse_outline_str_still_works():
    """str 输入路径不回退（原行为保持）。"""
    pms = _import_pptx_mcp_server()

    slides = pms._parse_outline("# 标题\n- 要点一\n- 要点二")
    assert len(slides) == 1
    assert slides[0]["title"] == "标题"
    assert len(slides[0]["points"]) == 2


def test_ppt_gen_from_list_outline(tmp_path):
    """list 大纲 → generate_ppt 真实生成 .pptx（备课→PPT 链路联通）。"""
    pms = _import_pptx_mcp_server()

    outline = [
        {"slide": 1, "title": "封面：函数的单调性", "points": ["数学 · 高中", "核心概念"]},
        {"slide": 2, "title": "定义", "points": ["严格定义", "几何意义"]},
        {"slide": 3, "title": "判定", "points": ["定义法", "导数法"]},
    ]
    r = pms.generate_ppt("函数的单调性", outline=outline,
                         out_name=f"round3_probe_{os.getpid()}")
    assert r.get("ok"), f"PPT 生成失败: {r.get('error')}"
    path = r.get("path")
    assert path and os.path.exists(path), f"PPT 未落盘: {path}"
    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) >= 3, f"应至少 3 页，实际 {len(prs.slides)}"
